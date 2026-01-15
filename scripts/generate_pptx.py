#!/usr/bin/env python3
"""Generate the portfolio PPTX using python-pptx.

Requires: python-pptx (see requirements.txt)
"""

from pathlib import Path

from pptx import Presentation

SLIDES = [
    ("Fail-Aware Manipulation", "Pick & Place with failure recovery"),
    ("Problem", "Pick & place failures reduce throughput and safety."),
    ("System Architecture", "Gazebo + MoveIt2 + Task Manager + Metrics"),
    ("FSM Design", "Recovery-driven state machine with retries."),
    ("Failure/Recovery Scenarios", "IK, Collision, Timeout, Goal not reached."),
    ("Experiment Setup", "Presets × repeats, JSONL logging."),
    ("Results", "Success rate, planning time, recovery rate."),
    ("Demo", "Insert Gazebo/RViz snapshots + video links."),
    ("Takeaways", "Structured recovery improves robustness."),
]


def main() -> None:
    prs = Presentation()
    for title, body in SLIDES:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body

    out_path = Path("ppt/Fail-Aware_Manipulation_Portfolio.pptx")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


if __name__ == "__main__":
    main()
